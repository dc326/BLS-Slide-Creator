# -*- coding: utf-8 -*-
"""
Created on Sun Jun  2 19:54:36 2024

@author: Dan Cunningham

The Following Code is intended to perform the following tasks:

1. Pull data from the Bureau of Labor Statistics
2. Manipulate the Output to create specific single or trended data outputs
3. Place the content into a PowerPoint Slide using a preformated slide template and the inserting content in specific positions

Written by Daniel Cunningham. Please feel free to reach out with any questions.
DCunningham326@gmail.com

"""

import pandas as pd
import requests
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Simplified Data Extraction for Latest Value
def find_latest_value(jsoninput):
    latest_data = None
    for item in jsoninput['data']:
        if item.get('latest') == 'true':
            latest_data = item
            break

    return (
        jsoninput['seriesID'],
        latest_data['value'],
        latest_data['periodName'] +
        "-" +
        latest_data['year'])

# Function for Pulling Specific Trended Data from Program
def find_trended_data(jsoninput, data_dict_key, key):
    for item in jsoninput:
        if data_dict_key[item['seriesID']] == key:
            return item['data']


def format_tbox(tframeshape):
    # Center and Middle align the text
    for para in tframeshape.paragraphs:
        para.alignment = PP_ALIGN.CENTER  # Center alignment
        para.vertical_ancor = MSO_ANCHOR.MIDDLE  # Middle alignment

    # Set font size to 60 and allow shrink to fit
    run = tframeshape.paragraphs[0].runs[0]
    run.font.size = Pt(60)
    tframeshape.auto_size = True


# pull BLS data
headers = {'Content-type': 'application/json'}
# LNS140 is UNemployment Rate, LNS130 is Unemployment level, JTS is Labor
# force Job Openings, CES05 is Avg Hourly Earnings
data = json.dumps({"seriesid": ['LNS13000000',
                                'JTS000000000000000JOL',
                                'CES0500000003',
                                'LNS14000000'],
                   "startyear": "2015",
                   "endyear": "2024"})
p = requests.post(
    '	https://api.bls.gov/publicAPI/v2/timeseries/data/',
    data=data,
    headers=headers)
json_data = json.loads(p.text)

# Key for Adding Specific Text Designations to BLS Model Data
data_dict_key = {
    'LNS13000000': "Unemployment Level",
    'JTS000000000000000JOL': "Labor Force",
    'CES0500000003': "Avg Hourly Earnings",
    'LNS14000000': "Unemployment Rate"}

# Extract and arrange core data
Lvals_dict = {}

for each in json_data['Results']['series']:
    catch = find_latest_value(each)
    # print(catch)
    Lvals_dict[data_dict_key[catch[0]]] = (catch[1], catch[2])


# Extract Unemployment Rate Data
URT_Data = Lvals_dict["Unemployment Rate"]

# Extract Average Hourly Earnings Data
AHE_Data = Lvals_dict["Avg Hourly Earnings"]

# Develop Dataframe with Labor Force and Job Openings in the same array
UEL_Data = find_trended_data(
    json_data['Results']['series'],
    data_dict_key,
    "Unemployment Level")
LF_Data = find_trended_data(
    json_data['Results']['series'],
    data_dict_key,
    "Labor Force")

UEL_DF = pd.DataFrame(UEL_Data).drop(columns=['footnotes', 'latest'])
LF_DF = pd.DataFrame(LF_Data).drop(columns=['footnotes', 'latest'])

merged_df = pd.merge(
    UEL_DF, LF_DF, on=[
        'year', 'period', 'periodName'], suffixes=(
            '_UEL', '_LB'))

# Create mapping for 'date' sorting
month_map = {
    'January': 1, 'February': 2, 'March': 3, 'April': 4,
    'May': 5, 'June': 6, 'July': 7, 'August': 8,
    'September': 9, 'October': 10, 'November': 11, 'December': 12
}

# Create a new column with the month numbers
merged_df['month'] = merged_df['periodName'].map(month_map)

# Create a new datetime column
merged_df['date'] = pd.to_datetime(merged_df[['year', 'month']].assign(day=1))

merged_df['pdnamelist'] = merged_df['periodName'] + "-" + merged_df['year']

merged_df.sort_values(ascending=True, by=['date'], inplace=True)


# Create Powerpoint Presentation

# Powerpoint Template
template_path = 'pptx_tests/BLS Template Presentation.pptx'

# Import the Powerpoint Template from Path
prs = Presentation(template_path)

# Select the Layout
layout_index = 2  # Replace with the identified index of the desired layout
slide_layout = prs.slide_layouts[layout_index]

# add slide to the presentation
slide = prs.slides.add_slide(slide_layout)

# add UR Shape to Slide
URTbox = slide.shapes.add_textbox(
    Inches(1.25),
    Inches(2.2),
    Inches(2.6),
    Inches(1.0))
URTbox.text_frame.text = URT_Data[0] + '%'
format_tbox(URTbox.text_frame)

# add Avg Hourly Earnings to Slide
AHEbox = slide.shapes.add_textbox(
    Inches(1.25),
    Inches(4.45),
    Inches(2.6),
    Inches(1.0))
AHEbox.text_frame.text = '$' + AHE_Data[0]
format_tbox(AHEbox.text_frame)

# Add chart to Slide
chart_data = CategoryChartData()
chart_data.categories = merged_df['pdnamelist']
chart_data.add_series('Unemployment Level', merged_df['value_UEL'])
chart_data.add_series('Labor Force', merged_df['value_LB'])

x, y, cx, cy = Inches(5.4), Inches(1.55), Inches(7), Inches(4.4)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# Format the chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# Set legend font size
chart.legend.font.size = Pt(12)

# Set axes font sizes
category_axis = chart.category_axis
value_axis = chart.value_axis
category_axis.tick_labels.font.size = Pt(11)
value_axis.tick_labels.font.size = Pt(11)

# Set vertical axis number format
value_axis.number_format = '#,##0'

# Add title to the chart
start_date = merged_df['date'].min().strftime('%b %Y')
end_date = merged_df['date'].max().strftime('%b %Y')
chart_title = chart.chart_title
chart_title.has_text_frame = True
chart_title.text_frame.text = f"Unemployment Level vs Labor Force\n{start_date} - {end_date}"
chart_title.text_frame.paragraphs[0].font.size = Pt(14)

# Add bottom textbox to slide with stock language "[Your Input Here]"
stock_txtbox = slide.shapes.add_textbox(
    Inches(0.81), Inches(6.21), Inches(11.72), Inches(1.0))
stock_txtbox.text_frame.text = '::Your Commentary Here::'


# Add Notes to the slide: Latest Data Available in BLS Data for All Datatables.
notes_slide = slide.notes_slide
notes_text_frame = notes_slide.notes_text_frame
notes_text_frame.text = "Unemployment Rate Data: " + \
    URT_Data[1] + "\n" + "Avg Hourly Earnings Data: " + AHE_Data[1]


# Save Presentation
prs.save('pptx_tests/Presentation_Template-test1.pptx')
